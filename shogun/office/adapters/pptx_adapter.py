"""PowerPoint Adapter — Hybrid python-pptx + COM automation for presentations.

Uses python-pptx for all slide manipulation (cross-platform, no Office needed).
Uses COM (via com_thread_pool) only for PDF export.

Placeholder convention: ``{{key}}`` — standard mustache style.
"""

from __future__ import annotations

import logging
import time
from pathlib import Path
from typing import Any

log = logging.getLogger("shogun.office.adapters.pptx")


# ── Presentation Handle ─────────────────────────────────────────────


class PresentationHandle:
    """Tracks an open PowerPoint presentation (python-pptx-based)."""

    def __init__(self, path: Path, presentation: Any):
        self.path = path
        self.presentation = presentation
        self.opened_at = time.time()


# ── Core Adapter Functions ───────────────────────────────────────────


def open_presentation(file_path: str) -> PresentationHandle:
    """Open a PowerPoint presentation with python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PowerPoint operations. Install with: pip install python-pptx")

    path = Path(file_path)
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        error_str = str(exc).lower()
        if "password" in error_str or "encrypted" in error_str:
            from shogun.office.exceptions import PasswordProtectedError
            raise PasswordProtectedError(str(path))
        if "corrupt" in error_str or "not a zip" in error_str or "package not found" in error_str:
            from shogun.office.exceptions import CorruptedFileError
            raise CorruptedFileError(str(path), str(exc))
        raise

    log.debug("Opened presentation: %s (%d slides)", path, len(prs.slides))
    return PresentationHandle(path=path, presentation=prs)


def close_presentation(handle: PresentationHandle) -> None:
    """Close a presentation (cleanup reference)."""
    handle.presentation = None
    log.debug("Closed presentation: %s", handle.path)


def list_slides(handle: PresentationHandle) -> list[dict[str, Any]]:
    """List all slides with their index and title."""
    prs = handle.presentation
    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type is not None:
                # Try to find the title shape
                if hasattr(shape, 'name') and 'title' in shape.name.lower():
                    title = shape.text_frame.text
                    break
        if not title:
            # Fallback: use first text shape
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text[:100]
                    break
        slides.append({
            "index": i,
            "title": title,
            "shape_count": len(slide.shapes),
        })
    return slides


def read_slide_text(handle: PresentationHandle, slide_index: int) -> str:
    """Read all text from a specific slide."""
    prs = handle.presentation
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (max {len(prs.slides) - 1}).")

    slide = prs.slides[slide_index]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return "\n".join(texts)


def replace_text(
    handle: PresentationHandle,
    slide_index: int,
    old_text: str,
    new_text: str,
) -> int:
    """Replace text on a specific slide.

    Returns the number of replacements made.
    """
    prs = handle.presentation
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range.")

    slide = prs.slides[slide_index]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        count += 1
    return count


def replace_placeholders(
    handle: PresentationHandle,
    mapping: dict[str, str],
) -> dict[str, int]:
    """Replace {{placeholder}} patterns across all slides.

    Args:
        mapping: Dict of placeholder → replacement value.

    Returns:
        Dict of placeholder → replacement count.
    """
    prs = handle.presentation
    counts: dict[str, int] = {k: 0 for k in mapping}

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in mapping.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)
                                counts[key] += 1
            # Also check table cells
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for key, value in mapping.items():
                                    if key in run.text:
                                        run.text = run.text.replace(key, value)
                                        counts[key] += 1

    log.debug("Replaced placeholders: %s", {k: v for k, v in counts.items() if v > 0})
    return counts


def insert_textbox(
    handle: PresentationHandle,
    slide_index: int,
    text: str,
    left: int = 1000000,
    top: int = 1000000,
    width: int = 5000000,
    height: int = 1000000,
) -> None:
    """Insert a text box on a slide.

    Dimensions are in EMUs (English Metric Units). 1 inch = 914400 EMU.
    """
    from pptx.util import Emu

    prs = handle.presentation
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range.")

    slide = prs.slides[slide_index]
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.text = text
    log.debug("Inserted textbox on slide %d", slide_index)


def insert_table(
    handle: PresentationHandle,
    slide_index: int,
    headers: list[str],
    rows: list[list[Any]],
    left: int = 500000,
    top: int = 2000000,
    width: int = 8000000,
    height: int = 1500000,
) -> None:
    """Insert a table on a slide."""
    from pptx.util import Emu

    prs = handle.presentation
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range.")

    slide = prs.slides[slide_index]
    total_rows = 1 + len(rows)
    total_cols = len(headers)

    table_shape = slide.shapes.add_table(
        total_rows, total_cols,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    table = table_shape.table

    # Headers
    for i, header in enumerate(headers):
        table.cell(0, i).text = str(header)

    # Data
    for r, row_data in enumerate(rows):
        for c, value in enumerate(row_data):
            if c < total_cols:
                table.cell(r + 1, c).text = str(value) if value is not None else ""

    log.debug("Inserted table on slide %d: %d×%d", slide_index, total_cols, len(rows))


def insert_image(
    handle: PresentationHandle,
    slide_index: int,
    image_path: str,
    left: int = 1000000,
    top: int = 2000000,
    width: int | None = None,
    height: int | None = None,
) -> None:
    """Insert an image on a slide.

    The image_path must be from an approved folder (validated by caller).
    """
    from pptx.util import Emu

    prs = handle.presentation
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range.")

    img_path = Path(image_path)
    if not img_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    slide = prs.slides[slide_index]
    kwargs: dict[str, Any] = {
        "image_file": str(img_path),
        "left": Emu(left),
        "top": Emu(top),
    }
    if width:
        kwargs["width"] = Emu(width)
    if height:
        kwargs["height"] = Emu(height)

    slide.shapes.add_picture(**kwargs)
    log.debug("Inserted image on slide %d: %s", slide_index, img_path.name)


def duplicate_slide(handle: PresentationHandle, slide_index: int) -> int:
    """Duplicate a slide. Returns the index of the new slide.

    Note: python-pptx has limited slide duplication support.
    This creates a copy of the slide layout and content.
    """
    import copy
    from lxml import etree

    prs = handle.presentation
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range.")

    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy shapes from template to new slide
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

    new_index = len(prs.slides) - 1
    log.debug("Duplicated slide %d → new slide %d", slide_index, new_index)
    return new_index


def save_as(handle: PresentationHandle, output_path: str) -> str:
    """Save the presentation to a new file path."""
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    handle.presentation.save(str(out))
    log.info("Saved presentation to: %s", out)
    return str(out)


def get_presentation_metadata(handle: PresentationHandle) -> dict[str, Any]:
    """Get metadata about the presentation."""
    prs = handle.presentation
    props = prs.core_properties
    return {
        "file": str(handle.path),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "title": props.title or "",
        "author": props.author or "",
        "created": str(props.created) if props.created else "",
        "modified": str(props.modified) if props.modified else "",
        "last_modified_by": props.last_modified_by or "",
    }


# ── COM-Only Functions (PDF Export) ──────────────────────────────────


def _com_export_pdf(file_path: str, output_path: str, visible: bool = False) -> str:
    """Export a presentation to PDF via COM.

    Must be called on the STA thread pool.
    """
    import win32com.client
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    # PowerPoint must be visible to open files
    ppt.Visible = True
    out = str(Path(output_path).resolve())
    try:
        pres = ppt.Presentations.Open(
            str(Path(file_path).resolve()),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False if not visible else True,
        )
        try:
            # ppSaveAsPDF = 32
            pres.SaveAs(out, FileFormat=32)
        finally:
            pres.Close()
    finally:
        ppt.Quit()

    log.info("Exported presentation to PDF: %s", out)
    return out


async def export_pdf(file_path: str, output_path: str, visible: bool = False) -> str:
    """Export to PDF — async wrapper around COM call."""
    from shogun.office.com_thread_pool import run_com, office_lock
    async with office_lock("powerpoint"):
        return await run_com(_com_export_pdf, file_path, output_path, visible)
